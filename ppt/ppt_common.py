# Este fichero contiene utilidades comunes para construir presentaciones PowerPoint a partir de gráficos y plantillas.
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt


def get_layout_by_name(prs: Presentation, layout_name: str):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    raise ValueError(f'No se ha encontrado el layout "{layout_name}" en la plantilla.')


def set_title_placeholder(slide, text: str) -> None:
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
            shape.text = text
            tf = shape.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].font.size = Pt(24)
            return


def set_footer_placeholder(slide, text: str) -> None:
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        current = shape.text.strip() if hasattr(shape, 'text') else ''
        if 'IndraMind' in current or 'dd/mm/aaaa' in current:
            shape.text = text
            return


def add_cover_slide(prs: Presentation, report_title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Portada Básica'))
    title_done = False
    subtitle_done = False

    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        ptype = shape.placeholder_format.type
        if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] and not title_done:
            shape.text = report_title
            title_done = True
        elif ptype in [PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY] and not subtitle_done:
            shape.text = subtitle
            subtitle_done = True


def add_section_slide(prs: Presentation, section_code: str, section_title: str, footer_text: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Separata secundaria'))

    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        text = shape.text.strip() if hasattr(shape, 'text') else ''
        ptype = shape.placeholder_format.type
        if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
            shape.text = section_title
        elif text == '2.1':
            shape.text = section_code
        elif 'IndraMind' in text or 'dd/mm/aaaa' in text:
            shape.text = footer_text


def add_chart_slide(prs: Presentation, title: str, image_path: Path, footer_text: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Diapositiva base Azul Amazónico'))
    set_title_placeholder(slide, title)
    set_footer_placeholder(slide, footer_text)
    slide.shapes.add_picture(str(image_path), Inches(0.65), Inches(1.75), Inches(12.0), Inches(4.95))