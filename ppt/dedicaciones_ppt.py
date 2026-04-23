# Este fichero genera la presentación PowerPoint de comité de dirección para el módulo de dedicaciones a partir del dataset filtrado.
from datetime import datetime
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory

import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt

from ppt.ppt_common import add_chart_slide, add_cover_slide, add_section_slide, get_layout_by_name, set_footer_placeholder, set_title_placeholder
from ui_common import DISPLAY_COLUMNS, PALETTE, PLOTLY_COLOR_SEQUENCE, format_number


def aggregate_monthly_totals(df: pd.DataFrame) -> pd.DataFrame:
    monthly = df.groupby('periodo', dropna=False, as_index=False).agg(horas_aplicadas=('horas_aplicadas', 'sum'), cantidad=('cantidad', 'sum')).sort_values('periodo')
    return monthly


def build_monthly_overview_figure(df_plot: pd.DataFrame, title: str):
    monthly_long = df_plot.melt(id_vars=['periodo'], value_vars=['horas_aplicadas', 'cantidad'], var_name='metrica', value_name='valor')
    monthly_long['metrica'] = monthly_long['metrica'].replace({'horas_aplicadas': 'Horas', 'cantidad': 'Cantidad (€)'})
    fig = px.bar(monthly_long, x='periodo', y='valor', color='metrica', barmode='group', title=title, color_discrete_map={'Horas': PALETTE['turquesa'], 'Cantidad (€)': '#5E7F8A'})
    fig.update_traces(hovertemplate='%{x}<br>%{fullData.name}: %{y:,.2f}<extra></extra>')
    fig.update_layout(height=460, paper_bgcolor='white', plot_bgcolor='white', font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'], size=22), xaxis=dict(title=dict(text='Periodo', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), yaxis=dict(title=dict(text='Valor', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='left', x=0, font=dict(color=PALETTE['texto_oscuro']), title_text=''), margin=dict(l=10, r=10, t=60, b=10))
    return fig


def aggregate_for_dimension(df: pd.DataFrame, dimension: str, metric: str) -> pd.DataFrame:
    grouped = df.groupby(dimension, dropna=False, as_index=False).agg(horas_aplicadas=('horas_aplicadas', 'sum'), cantidad=('cantidad', 'sum'), empleados=('empleado', 'nunique'), departamentos=('departamento', 'nunique'))
    grouped = grouped.sort_values(metric, ascending=False)
    return grouped


def aggregate_timeline(df: pd.DataFrame, dimension: str, metric: str) -> pd.DataFrame:
    return df.groupby(['periodo', dimension], dropna=False, as_index=False)[metric].sum().sort_values(['periodo', metric], ascending=[True, False])


def build_bar_figure(df_plot: pd.DataFrame, x_col: str, y_col: str, title: str):
    plot_df = df_plot.copy()
    fig = px.bar(plot_df, x=y_col, y=x_col, orientation='h', title=title, color=x_col, color_discrete_sequence=PLOTLY_COLOR_SEQUENCE)
    fig.update_traces(marker_line_color=PALETTE['gris_ceramica'], marker_line_width=0.4, hovertemplate='%{y}<br>%{x:,.2f}<extra></extra>')
    fig.update_layout(height=max(520, 28 * len(plot_df) + 120), paper_bgcolor='white', plot_bgcolor='white', font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'], size=22), xaxis=dict(title=dict(text=DISPLAY_COLUMNS.get(y_col, y_col), font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), yaxis=dict(title=dict(text='', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.06)'), showlegend=False, margin=dict(l=10, r=10, t=50, b=10))
    return fig


def build_timeline_figure(df_plot: pd.DataFrame, dimension: str, metric: str, title: str):
    fig = px.line(df_plot, x='periodo', y=metric, color=dimension, markers=True, title=title, color_discrete_sequence=PLOTLY_COLOR_SEQUENCE)

    for idx, trace in enumerate(fig.data):
        color = PLOTLY_COLOR_SEQUENCE[idx % len(PLOTLY_COLOR_SEQUENCE)]
        trace.line.width = 2
        trace.marker.size = 5
        trace.line.color = color
        trace.marker.color = color

    fig.update_layout(height=460, paper_bgcolor='white', plot_bgcolor='white', font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'], size=22), xaxis=dict(title=dict(text='Periodo', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), yaxis=dict(title=dict(text=DISPLAY_COLUMNS.get(metric, metric), font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), legend=dict(orientation='v', yanchor='top', y=1, xanchor='left', x=1.02, font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'])), margin=dict(l=10, r=220, t=50, b=10))
    return fig


def export_figure_png(fig, output_path: Path) -> Path:
    fig.write_image(str(output_path), format='png', scale=2, width=1600, height=900)
    return output_path


def add_summary_slide(prs: Presentation, metrics: dict, footer_text: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Diapositiva base Azul Amazónico'))
    set_title_placeholder(slide, 'Resumen ejecutivo del proyecto')
    set_footer_placeholder(slide, footer_text)

    box_w = 2.6
    box_h = 1.2
    x_positions = [0.8, 3.55, 6.3, 9.05]
    y = 2.15

    for idx, item in enumerate(metrics.items()):
        label, value = item
        shape = slide.shapes.add_shape(1, Inches(x_positions[idx]), Inches(y), Inches(box_w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string('E3E2DA')
        shape.line.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('004254')

        text_frame = shape.text_frame
        text_frame.clear()
        p1 = text_frame.paragraphs[0]
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('004254')

        p2 = text_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = value
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('004254')

    body = slide.shapes.add_textbox(Inches(0.9), Inches(4.15), Inches(11.2), Inches(1.3))
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Presentación de comité de dirección elaborada con el dataset completo del proyecto. Incluye vistas por departamento/elemento y por empleado, tanto en horas como en cantidad.'
    r.font.size = Pt(18)
    r.font.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('E3E2DA')


def build_committee_presentation(df: pd.DataFrame, template_path: str, report_title: str = 'Informe de Costes del Proyecto', document_name: str = 'Informe de Costes del Proyecto') -> bytes:
    prs = Presentation(template_path)
    footer_text = f'IndraMind • {document_name} • {datetime.now().strftime("%d/%m/%Y")}'

    with TemporaryDirectory(prefix='indra_ppt_') as temp_dir_str:
        temp_dir = Path(temp_dir_str)

        add_cover_slide(prs, report_title, 'Presentación ejecutiva de comité de dirección')

        summary_metrics = {
            'Coste total': format_number(df['cantidad'].sum(), 2),
            'Horas totales': format_number(df['horas_aplicadas'].sum(), 2),
            'Departamentos': format_number(df['departamento'].nunique(), 0),
            'Empleados': format_number(df['empleado'].nunique(), 0)
        }
        add_summary_slide(prs, summary_metrics, footer_text)

        monthly_totals = aggregate_monthly_totals(df)
        monthly_overview_fig = build_monthly_overview_figure(monthly_totals, 'Evolución mensual global · Horas y Cantidad (€)')
        monthly_overview_path = export_figure_png(monthly_overview_fig, temp_dir / '00_resumen_mensual_global.png')
        add_chart_slide(prs, 'Resumen general · Evolución mensual de horas y cantidad (€)', monthly_overview_path, footer_text)

        sections = [
            ('1.1', 'Elemento + Departamento (Horas)', 'departamento', 'horas_aplicadas'),
            ('1.2', 'Empleado + Nombre (Horas)', 'empleado', 'horas_aplicadas'),
            ('1.3', 'Elemento + Departamento (Cantidad)', 'departamento', 'cantidad'),
            ('1.4', 'Empleado + Nombre (Cantidad)', 'empleado', 'cantidad')
        ]

        for section_code, section_title, dimension, metric in sections:
            add_section_slide(prs, section_code, section_title, footer_text)
            aggregated = aggregate_for_dimension(df, dimension, metric)
            timeline = aggregate_timeline(df, dimension, metric)
            ranking_fig = build_bar_figure(aggregated, dimension, metric, f'{section_title} - ranking general')
            ranking_path = export_figure_png(ranking_fig, temp_dir / f'{section_code.replace(".", "_")}_ranking.png')
            add_chart_slide(prs, f'{section_title} · Ranking general', ranking_path, footer_text)
            evolution_fig = build_timeline_figure(timeline, dimension, metric, f'{section_title} - evolución mensual')
            evolution_path = export_figure_png(evolution_fig, temp_dir / f'{section_code.replace(".", "_")}_evolucion.png')
            add_chart_slide(prs, f'{section_title} · Evolución mensual', evolution_path, footer_text)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()