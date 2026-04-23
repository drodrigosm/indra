# Este fichero genera una presentación PowerPoint de comité de dirección a partir del dataset procesado en Streamlit, reutilizando una plantilla PPTX y exportando slides de resumen y gráficos.
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from datetime import datetime

import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER


PALETTE = {
    'azul_oscuro': '#001923',
    'turquesa': '#00B0BD',
    'gris_ceramica': '#E3E2DA',
    'azul_amazonico': '#004254',
    'texto_claro': '#E3E2DA',
    'texto_oscuro': '#001923'
}

PLOTLY_COLOR_SEQUENCE = ['#00B0BD', '#33C2CD', '#66D0D6', '#0097A7', '#26AAB5', '#4DBCC4', '#7BCFD4', '#5E7F8A', '#6F909A', '#80A1AA', '#91B2BA', '#A3C2C9', '#8FA3AA', '#9FB2B8', '#AFC1C6', '#BFCFD3', '#C9D4D6', '#B8C5C8', '#A7B6BA', '#96A8AD']

DISPLAY_COLUMNS = {
    'departamento': 'Departamento',
    'empleado': 'Empleado',
    'elemento': 'Elemento',
    'nombre': 'Nombre',
    'fecha': 'Periodo',
    'periodo': 'Periodo',
    'horas_aplicadas': 'Horas Aplicadas',
    'cantidad': 'Cantidad',
    'tasa': 'Tasa',
    'categoria_nombre': 'Categoría',
    'tipo_coste_nombre': 'Tipo de Coste'
}

def aggregate_monthly_totals(df: pd.DataFrame) -> pd.DataFrame:
    monthly = df.groupby('periodo', dropna=False, as_index=False).agg(horas_aplicadas=('horas_aplicadas', 'sum'), cantidad=('cantidad', 'sum')).sort_values('periodo')
    return monthly

def build_monthly_overview_figure(df_plot: pd.DataFrame, title: str):
    monthly_long = df_plot.melt(id_vars=['periodo'], value_vars=['horas_aplicadas', 'cantidad'], var_name='metrica', value_name='valor')
    monthly_long['metrica'] = monthly_long['metrica'].replace({'horas_aplicadas': 'Horas', 'cantidad': 'Cantidad (€)'})
    fig = px.bar(monthly_long, x='periodo', y='valor', color='metrica', barmode='group', title=title, color_discrete_map={'Horas': PALETTE['turquesa'], 'Cantidad (€)': '#5E7F8A'})
    fig.update_traces(hovertemplate='%{x}<br>%{fullData.name}: %{y:,.2f}<extra></extra>')
    fig.update_layout(height=460, paper_bgcolor='white', plot_bgcolor='white', font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'], size=22), xaxis=dict(title=dict(text='Periodo', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), yaxis=dict(title=dict(text='Valor', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='left', x=0, font=dict(color=PALETTE['texto_oscuro']), title_text=''), margin=dict(l=10, r=10, t=60, b=10))
    return fig

def format_number(value: float, decimals: int = 2) -> str:
    return f"{value:,.{decimals}f}".replace(',', 'X').replace('.', ',').replace('X', '.')


def aggregate_for_dimension(df: pd.DataFrame, dimension: str, metric: str) -> pd.DataFrame:
    grouped = df.groupby(dimension, dropna=False, as_index=False).agg(horas_aplicadas=('horas_aplicadas', 'sum'), cantidad=('cantidad', 'sum'), empleados=('empleado', 'nunique'), departamentos=('departamento', 'nunique'))
    grouped = grouped.sort_values(metric, ascending=False)
    return grouped


def aggregate_timeline(df: pd.DataFrame, dimension: str, metric: str) -> pd.DataFrame:
    return df.groupby(['periodo', dimension], dropna=False, as_index=False)[metric].sum().sort_values(['periodo', metric], ascending=[True, False])


def build_bar_figure(df_plot: pd.DataFrame, x_col: str, y_col: str, title: str):
    plot_df = df_plot.copy()
    fig = px.bar(plot_df, x=y_col, y=x_col, orientation='h', title=title, color=x_col, color_discrete_sequence=PLOTLY_COLOR_SEQUENCE)
    fig.update_traces(marker_line_color=PALETTE['gris_ceramica'], marker_line_width=0.4, hovertemplate='%{y}<br>%{x:,.2f}<extra></extra>')
    fig.update_layout(height=max(520, 28 * len(plot_df) + 120), paper_bgcolor='white', plot_bgcolor='white', font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'], size=22), xaxis=dict(title=dict(text=DISPLAY_COLUMNS.get(y_col, y_col), font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), yaxis=dict(title=dict(text='', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.06)'), showlegend=False, margin=dict(l=10, r=10, t=50, b=10))
    return fig


def build_timeline_figure(df_plot: pd.DataFrame, dimension: str, metric: str, title: str):
    fig = px.line(df_plot, x='periodo', y=metric, color=dimension, markers=True, title=title, color_discrete_sequence=PLOTLY_COLOR_SEQUENCE)
    for idx, trace in enumerate(fig.data):
        color = PLOTLY_COLOR_SEQUENCE[idx % len(PLOTLY_COLOR_SEQUENCE)]
        trace.line.width = 2
        trace.marker.size = 5
        trace.line.color = color
        trace.marker.color = color
    fig.update_layout(height=460, paper_bgcolor='white', plot_bgcolor='white', font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'], size=22), xaxis=dict(title=dict(text='Periodo', font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), yaxis=dict(title=dict(text=DISPLAY_COLUMNS.get(metric, metric), font=dict(color=PALETTE['texto_oscuro'])), tickfont=dict(color=PALETTE['texto_oscuro']), gridcolor='rgba(0,66,84,0.10)'), legend=dict(orientation='v', yanchor='top', y=1, xanchor='left', x=1.02, font=dict(color=PALETTE['texto_oscuro']), title_font=dict(color=PALETTE['azul_amazonico'])), margin=dict(l=10, r=220, t=50, b=10))
    return fig


def export_figure_png(fig, output_path: Path) -> Path:
    fig.write_image(str(output_path), format='png', scale=2, width=1600, height=900)
    return output_path


def get_layout_by_name(prs: Presentation, layout_name: str):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    raise ValueError(f'No se ha encontrado el layout "{layout_name}" en la plantilla.')


def set_first_text_placeholder(slide, text: str) -> None:
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY]:
            shape.text = text
            return


def set_title_placeholder(slide, text: str) -> None:
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
            shape.text = text
            tf = shape.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].font.size = Pt(24)
            return


def set_footer_placeholder(slide, text: str) -> None:
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        current = shape.text.strip() if hasattr(shape, 'text') else ''
        if 'IndraMind' in current or 'dd/mm/aaaa' in current:
            shape.text = text
            return


def add_cover_slide(prs: Presentation, report_title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Portada Básica'))
    title_done = False
    subtitle_done = False
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        ptype = shape.placeholder_format.type
        if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] and not title_done:
            shape.text = report_title
            title_done = True
        elif ptype in [PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY] and not subtitle_done:
            shape.text = subtitle
            subtitle_done = True


def add_section_slide(prs: Presentation, section_code: str, section_title: str, footer_text: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Separata secundaria'))
    for shape in slide.placeholders:
        if not hasattr(shape, 'text_frame'):
            continue
        text = shape.text.strip() if hasattr(shape, 'text') else ''
        ptype = shape.placeholder_format.type
        if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
            shape.text = section_title
        elif text == '2.1':
            shape.text = section_code
        elif 'IndraMind' in text or 'dd/mm/aaaa' in text:
            shape.text = footer_text


def add_summary_slide(prs: Presentation, metrics: dict, footer_text: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Diapositiva base Azul Amazónico'))
    set_title_placeholder(slide, 'Resumen ejecutivo del proyecto')
    set_footer_placeholder(slide, footer_text)

    box_w = 2.6
    box_h = 1.2
    x_positions = [0.8, 3.55, 6.3, 9.05]
    y = 2.15

    for idx, item in enumerate(metrics.items()):
        label, value = item
        shape = slide.shapes.add_shape(1, Inches(x_positions[idx]), Inches(y), Inches(box_w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string('E3E2DA')
        shape.line.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('004254')

        text_frame = shape.text_frame
        text_frame.clear()
        p1 = text_frame.paragraphs[0]
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('004254')

        p2 = text_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = value
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('004254')

    body = slide.shapes.add_textbox(Inches(0.9), Inches(4.15), Inches(11.2), Inches(1.3))
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Presentación de comité de dirección elaborada con el dataset completo del proyecto. Incluye vistas por departamento/elemento y por empleado, tanto en horas como en cantidad.'
    r.font.size = Pt(18)
    r.font.color.rgb = __import__('pptx').dml.color.RGBColor.from_string('E3E2DA')


def add_chart_slide(prs: Presentation, title: str, image_path: Path, footer_text: str) -> None:
    slide = prs.slides.add_slide(get_layout_by_name(prs, 'Diapositiva base Azul Amazónico'))
    set_title_placeholder(slide, title)
    set_footer_placeholder(slide, footer_text)
    slide.shapes.add_picture(str(image_path), Inches(0.65), Inches(1.75), Inches(12.0), Inches(4.95))


def build_committee_presentation(df: pd.DataFrame, template_path: str, report_title: str = 'Informe de Costes del Proyecto', document_name: str = 'Informe de Costes del Proyecto') -> bytes:
    prs = Presentation(template_path)
    footer_text = f'IndraMind • {document_name} • {datetime.now().strftime("%d/%m/%Y")}'

    with TemporaryDirectory(prefix='indra_ppt_') as temp_dir_str:
        temp_dir = Path(temp_dir_str)

        add_cover_slide(prs, report_title, 'Presentación ejecutiva de comité de dirección')

        summary_metrics = {
            'Coste total': format_number(df['cantidad'].sum(), 2),
            'Horas totales': format_number(df['horas_aplicadas'].sum(), 2),
            'Departamentos': format_number(df['departamento'].nunique(), 0),
            'Empleados': format_number(df['empleado'].nunique(), 0)
        }
        add_summary_slide(prs, summary_metrics, footer_text)

        monthly_totals = aggregate_monthly_totals(df)
        monthly_overview_fig = build_monthly_overview_figure(monthly_totals, 'Evolución mensual global · Horas y Cantidad (€)')
        monthly_overview_path = export_figure_png(monthly_overview_fig, temp_dir / '00_resumen_mensual_global.png')
        add_chart_slide(prs, 'Resumen general · Evolución mensual de horas y cantidad (€)', monthly_overview_path, footer_text)

        sections = [
            ('1.1', 'Elemento + Departamento (Horas)', 'departamento', 'horas_aplicadas'),
            ('1.2', 'Empleado + Nombre (Horas)', 'empleado', 'horas_aplicadas'),
            ('1.3', 'Elemento + Departamento (Cantidad)', 'departamento', 'cantidad'),
            ('1.4', 'Empleado + Nombre (Cantidad)', 'empleado', 'cantidad')
        ]

        for section_code, section_title, dimension, metric in sections:
            add_section_slide(prs, section_code, section_title, footer_text)

            aggregated = aggregate_for_dimension(df, dimension, metric)
            timeline = aggregate_timeline(df, dimension, metric)

            ranking_fig = build_bar_figure(aggregated, dimension, metric, f'{section_title} - ranking general')
            ranking_path = export_figure_png(ranking_fig, temp_dir / f'{section_code.replace(".", "_")}_ranking.png')
            add_chart_slide(prs, f'{section_title} · Ranking general', ranking_path, footer_text)

            evolution_fig = build_timeline_figure(timeline, dimension, metric, f'{section_title} - evolución mensual')
            evolution_path = export_figure_png(evolution_fig, temp_dir / f'{section_code.replace(".", "_")}_evolucion.png')
            add_chart_slide(prs, f'{section_title} · Evolución mensual', evolution_path, footer_text)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()